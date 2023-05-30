from pptx import Presentation
from pptx.util import Pt

def create_title(prs, title, sub_title):
    # Add a title slide
    title_slide_layout = prs.slide_layouts[0]              #投影片封面
    slide1 = prs.slides.add_slide(title_slide_layout)      #新增投影片
    
    title1 = slide1.shapes.title                           #設定投影片標題
    title1.text = title
    
    subtitle = slide1.placeholders[1]                      #副標題
    subtitle.text = sub_title

def create_body(prs, slides):
    # Loop through each slide and extract the title and content fields
    for slide in slides:                                   #取每個投影片的標題與內容
        slide_title = slide['title']
        slide_content = slide['content']

        # Add a bullet slide
        bullet_slide_layout = prs.slide_layouts[1]         #第一個內容的投影片
        slide2 = prs.slides.add_slide(bullet_slide_layout) #新增投影片
        
        title2 = slide2.shapes.title                       #設定投影片標題
        title2.text = slide_title
        
        body2 = slide2.shapes.placeholders[1]              #主要文本的內容
        tf = body2.text_frame                              #矩形框

        tf.clear()
        
        for content_string in slide_content:
            p = tf.paragraphs[-1]                          #將每個新增的item加在最後一項
            
            # set font size
            run = p.add_run()                              
            run.text = content_string                      
            font = run.font                                
            font.name = '微軟正黑體'                          #字型
            font.size = Pt(18)                             #字體大小
            
            p.level = 0                                    #level=0為第一層
            p = tf.add_paragraph()

if __name__ == '__main__':
    prs = Presentation()                                   #建一個PPT
    
    title = '如何準備英文報告'
    subtitle = 'subtitle'
    create_title(prs, title, subtitle)
    
    slides = []
    slides.append({'title': '進行步驟一', 'content': ['item 1', 'item 2', 'item 3']})
    slides.append({'title': '進行步驟二', 'content': ['item 1', 'item 2']})
    create_body(prs, slides)
    
    prs.save('./1.pptx')    
    
